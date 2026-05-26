import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Paths to assets
brain_dir = Path(r"C:/Users/15438/.gemini/antigravity/brain/731adfd5-7f8d-4f4b-90bf-7a544a8b1640")
arch_image = brain_dir / "spring_boot_arch_diagram_mysql_1767629272144.png"
placeholder_image = brain_dir / "demo_screenshot_placeholder_1767625816386.png"

# Output PPT path
output_path = Path(__file__).parent / "College_Project_Presentation.pptx"

# Helper: Apply modern teal theme
def apply_modern_theme(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xAA, 0x00, 0x00)  # Red background
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Inter"
                if run.font.size is None:
                    run.font.size = Pt(20)

# Helper: Add Title
def add_title(slide, text):
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = text
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Inter"
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Helper: Add Bullet Points
def add_bullets(slide, points):
    # Use existing body placeholder if available, else create textbox
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
    else:
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
    
    tf.clear()
    for p_text in points:
        p = tf.add_paragraph()
        p.text = p_text
        p.level = 0
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.name = "Inter"
            run.font.size = Pt(24)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()

# Layouts
title_layout = prs.slide_layouts[0]
bullet_layout = prs.slide_layouts[1]
blank_layout = prs.slide_layouts[6]

# SLIDE 1: Project Title
slide = prs.slides.add_slide(title_layout)
slide.shapes.title.text = "Dynamic College Website CMS"
slide.placeholders[1].text = "Techno India Main Project\nSpring Boot Implementation"
apply_modern_theme(slide)

# SLIDE 2: Problem Statement
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Problem Statement"
add_bullets(slide, [
    "Static websites require manual code editing for every update.",
    "Lack of real-time communication between college admin and students.",
    "Difficulty in managing course curriculum and notices efficiently.",
    "Dependency on technical staff for minor content changes."
])
apply_modern_theme(slide)

# SLIDE 3: Objectives
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Objectives"
add_bullets(slide, [
    "To develop a dynamic Content Management System (CMS).",
    "To provide an intuitive Admin Interface for data management.",
    "To ensure real-time updates of academic programs and facilities.",
    "To implement a secure and scalable backend architecture."
])
apply_modern_theme(slide)

# SLIDE 4: Project Overview
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Project Overview"
add_bullets(slide, [
    "A full-stack web application replacing legacy static pages.",
    "Built with Java Spring Boot for robust backend processing.",
    "Features a decoupled architecture with Thymeleaf frontend.",
    "Focuses on usability, performance, and data integrity."
])
apply_modern_theme(slide)

# SLIDE 5: System Architecture (Diagram)
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "System Architecture")
slide.shapes.add_picture(str(arch_image), Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))
apply_modern_theme(slide)

# SLIDE 6: Architecture Explained
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Architecture Details"
add_bullets(slide, [
    "MVC Pattern: Model-View-Controller separation.",
    "Controller: REST endpoints receiving client requests.",
    "Service: Core business logic and transaction management.",
    "Repository: Data persistence using Spring Data JPA.",
    "Database: Relational schema storing structured data."
])
apply_modern_theme(slide)

# SLIDE 7: Key Features
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Key Features"
add_bullets(slide, [
    "Dynamic Content Management: Update courses/notices instantly.",
    "Admin Dashboard: Secure access to CRUD operations.",
    "Responsive UI: Mobile-friendly design with Bootstrap/CSS.",
    "Search & Filter: Easy navigation for academic programs.",
    "Scalable Database: Optimized for high data volume."
])
apply_modern_theme(slide)

# SLIDE 8: Technology Stack
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Technology Stack"
add_bullets(slide, [
    "Backend: Spring Boot 3.3.0, Java 17",
    "Frontend: Thymeleaf, HTML5, CSS3, JavaScript",
    "Database: MySQL / H2 In-Memory",
    "Tools: Maven, STS IDE, Git"
])
apply_modern_theme(slide)

# SLIDE 9: Database Design
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Database Design"
add_bullets(slide, [
    "Entities: Degree, Course, PageSection, PageContent.",
    "Relationships: One-to-Many mapping (e.g., Degree -> Courses).",
    "Normalization: ensuring data consistency and reducing redundancy.",
    "JPA Entities handling automatic table generation."
])
apply_modern_theme(slide)

# SLIDE 10: Demo Screenshots
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "System Demo")
slide.shapes.add_picture(str(placeholder_image), Inches(1), Inches(1.5), width=Inches(8), height=Inches(4.5))
apply_modern_theme(slide)

# SLIDE 11: Future Scope
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Future Scope"
add_bullets(slide, [
    "Integration with Payment Gateway for admissions.",
    "AI-based Chatbot for student queries.",
    "Role-based Access Control (RBAC) implementation.",
    "Mobile Application development using REST APIs."
])
apply_modern_theme(slide)

# SLIDE 12: Conclusion
slide = prs.slides.add_slide(bullet_layout)
slide.shapes.title.text = "Conclusion"
add_bullets(slide, [
    "Successfully delivered a functional Dynamic CMS.",
    "Streamlined the process of website content updates.",
    "Demonstrated effective use of Spring Boot ecosystem.",
    "Project is ready for deployment and future enhancements."
])
apply_modern_theme(slide)

# Save
prs.save(str(output_path))
print(f"Presentation generated at: {output_path}")
